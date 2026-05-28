from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


def apply_theme(slide, color=RGBColor(255, 255, 255)):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(0, 0, 0)


def add_title_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_theme(slide)
    slide.shapes.title.text = title
    if subtitle is not None:
        slide.placeholders[1].text = subtitle


def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_theme(slide)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.text = bullet
        p.level = 0


def add_flow_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    apply_theme(slide)
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.6)
    width = Inches(2.2)
    height = Inches(0.9)

    steps = [
        "Upload PDF",
        "Presign + B2 Upload",
        "Confirm + DB Row",
        "Ingest + Chunk",
        "Chroma Upsert",
        "Chat + Retrieve",
        "LLM Answer",
        "LLM Judge"
    ]

    shapes = []
    for i, label in enumerate(steps):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left + Inches((i % 4) * 2.4),
            top + Inches((i // 4) * 1.4),
            width,
            height,
        )
        shape.text_frame.text = label
        shapes.append(shape)


def build_deck(path):
    prs = Presentation()

    add_title_slide(
        prs,
        "RAG Threads Application",
        "Flask + Neon + Chroma Cloud + Backblaze B2 + OpenAI"
    )

    add_bullets_slide(prs, "Agenda", [
        "Problem & goal",
        "Architecture overview",
        "Core components",
        "Data flow",
        "Storage & vector search",
        "Chat + RAG pipeline",
        "Admin & security",
        "Evaluation metrics",
        "Testing & learnings"
    ])

    add_bullets_slide(prs, "Goal", [
        "Multi-tenant RAG app for PDF Q&A",
        "Global and user-private documents",
        "No general chat without a document",
        "Admin management and audit logs"
    ])

    add_bullets_slide(prs, "Tech Stack", [
        "Backend: Flask + SQLAlchemy + Alembic",
        "DB: Neon Postgres",
        "Vector DB: Chroma Cloud",
        "Storage: Backblaze B2 (S3 compatible)",
        "LLM: OpenAI (gpt-4o)",
        "Frontend: Server-rendered templates + JS"
    ])

    add_bullets_slide(prs, "Configuration Highlights", [
        "MAX_UPLOAD_MB = 100",
        "CHUNK_SIZE = 1000",
        "CHUNK_OVERLAP = 200",
        "TOP_K_CHUNKS = 5",
        "MAX_COMPLETION_TOKENS = 4000",
        "RAG_TEMPERATURE = 0.7"
    ])

    add_bullets_slide(prs, "Project Structure", [
        "app/models: ORM models",
        "app/routes: API + views",
        "app/services: auth, storage, RAG, evaluation",
        "app/templates + static: UI",
        "migrations: Alembic",
        "tests: unit & integration"
    ])

    add_bullets_slide(prs, "app/models Summary", [
        "user.py: accounts + roles",
        "document.py: PDF metadata + scope",
        "thread.py: conversation metadata",
        "message.py: user/assistant content",
        "message_evaluation.py: LLM-judge scores",
        "ingestion_job.py: processing queue",
        "activity_log.py: audit trail"
    ])

    add_bullets_slide(prs, "app/routes Summary", [
        "auth.py: register/login/OAuth",
        "documents.py: upload/presign/confirm",
        "threads.py: CRUD threads",
        "chat.py: RAG responses + metrics",
        "admin.py: admin-only management",
        "views.py: HTML pages"
    ])

    add_bullets_slide(prs, "app/services Summary", [
        "auth_service: sessions + guards",
        "storage_service: B2 presigned URLs",
        "vector_service: Chroma search/upsert",
        "rag_service: retrieval + prompt",
        "ingestion_service: PDF processing",
        "evaluation_service: LLM judge"
    ])

    add_bullets_slide(prs, "Key Models", [
        "users, documents, threads, messages",
        "ingestion_jobs for async processing",
        "activity_logs for audit",
        "message_evaluations for LLM-judge metrics"
    ])

    add_bullets_slide(prs, "Upload Flow", [
        "Client requests presigned URL",
        "Direct upload to B2",
        "Confirm upload creates document + ingestion job",
        "Worker extracts text, chunks, embeds"
    ])

    add_bullets_slide(prs, "Ingestion Worker", [
        "Poll queued ingestion jobs",
        "Download PDF from B2",
        "Extract text per page",
        "Chunk with overlap",
        "Embed with OpenAI",
        "Upsert to Chroma",
        "Update job status"
    ])

    add_bullets_slide(prs, "Ingestion & Chunking", [
        "Download PDF from B2",
        "Extract per-page text",
        "Chunk with overlap",
        "Generate embeddings",
        "Upsert into Chroma with metadata"
    ])

    add_bullets_slide(prs, "Chroma Cloud Usage", [
        "Single collection for vectors",
        "Metadata filters: document_id, scope, owner_user_id",
        "Search returns top-k chunks"
    ])

    add_bullets_slide(prs, "Backblaze B2 Usage", [
        "Store original PDFs",
        "Store extracted images (optional)",
        "Presigned upload/download",
        "Delete by prefix for cleanup"
    ])

    add_bullets_slide(prs, "RAG Chat Pipeline", [
        "Retrieve top-k chunks from Chroma",
        "Assemble prompt with excerpts",
        "Call OpenAI for answer",
        "Store response + citations"
    ])

    add_bullets_slide(prs, "Prompting & Temperature", [
        "System prompt enforces grounded answers",
        "RAG_TEMPERATURE = 0.7",
        "Max completion tokens = 4000",
        "Message history included for context"
    ])

    add_flow_slide(prs, "End-to-End Flow")

    add_bullets_slide(prs, "Authorization", [
        "login_required, admin_required",
        "owns_document, owns_thread",
        "Admin-only management endpoints"
    ])

    add_bullets_slide(prs, "Admin Panel", [
        "User management",
        "Document management",
        "Activity log",
        "Re-ingest / delete chunks"
    ])

    add_bullets_slide(prs, "Deletion Semantics", [
        "Soft-delete records",
        "Delete vectors by document_id",
        "Delete B2 objects by prefix",
        "Cascade on user deletion"
    ])

    add_bullets_slide(prs, "Evaluation (LLM Judge)", [
        "Faithfulness (0–1)",
        "Citation precision (0–1)",
        "Groundedness (0–1)",
        "Scores stored per message"
    ])

    add_bullets_slide(prs, "Evaluation Workflow", [
        "Extract citations from answer text",
        "Judge compares answer to excerpts",
        "Heuristics penalize missing citations",
        "Scores shown to users"
    ])

    add_bullets_slide(prs, "Evaluation UX", [
        "Metrics panel shows scores (0–100)",
        "Heuristics penalize missing citations",
        "Off-topic answers drop scores"
    ])

    add_bullets_slide(prs, "Testing", [
        "Pytest unit tests for models/auth",
        "Integration tests for access rules",
        "Evaluation persistence test"
    ])

    add_bullets_slide(prs, "Security Considerations", [
        "Role-based access control",
        "Scope filtering in vector search",
        "No cross-tenant data leakage"
    ])

    add_bullets_slide(prs, "Deletion Flow", [
        "Soft-delete records in Postgres",
        "Delete vectors by document_id",
        "Delete B2 objects by prefix",
        "Admin cascade delete users"
    ])

    add_bullets_slide(prs, "Performance Notes", [
        "Batch upserts to Chroma",
        "Streaming chat for UX",
        "Background ingestion worker"
    ])

    add_bullets_slide(prs, "Interview Talking Points", [
        "Why presigned uploads",
        "How metadata filters enforce access",
        "LLM judge for quality signals",
        "Operational tradeoffs"
    ])

    add_bullets_slide(prs, "Future Improvements", [
        "Better citation extraction",
        "Asynchronous judge jobs",
        "More evaluation metrics",
        "Team/tenant support"
    ])

    add_bullets_slide(prs, "Thank You", [
        "Questions?",
        "Demo ready"
    ])

    prs.save(path)


if __name__ == "__main__":
    build_deck("rag_deck.pptx")
